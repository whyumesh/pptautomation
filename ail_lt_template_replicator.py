"""
AIL LT PPT Template Replicator
Creates PPT exactly matching the manual template structure, formatting, and positioning
Uses the template as a base and populates with data from Excel files
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime
import logging
import argparse
import calendar
import json
import re

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class TemplateReplicator:
    """Replicates the template PPT exactly with data from Excel files"""
    
    def __init__(self, template_path: Path):
        """
        Initialize with template PPT.
        
        Args:
            template_path: Path to the manual template PPT
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        # Load template
        self.template_prs = Presentation(str(template_path))
        self.prs = Presentation(str(template_path))  # Start with template
        logger.info(f"Loaded template: {template_path}")
        
        # Load template analysis if available
        analysis_path = Path("template_analysis.json")
        if analysis_path.exists():
            with open(analysis_path, 'r', encoding='utf-8') as f:
                self.template_info = json.load(f)
        else:
            self.template_info = None
            logger.warning("Template analysis not found. Run analyze_template_complete.py first.")
    
    def read_excel_file(self, file_path: Path, sheet_name=None):
        """Read Excel file, handling both .xlsx and .xlsb formats"""
        try:
            if file_path.suffix == '.xlsb':
                try:
                    if sheet_name:
                        return pd.read_excel(file_path, sheet_name=sheet_name, engine='pyxlsb')
                    else:
                        return pd.read_excel(file_path, engine='pyxlsb')
                except ImportError:
                    logger.warning("pyxlsb not installed. Install with: pip install pyxlsb")
                    if sheet_name:
                        return pd.read_excel(file_path, sheet_name=sheet_name)
                    else:
                        return pd.read_excel(file_path)
            else:
                if sheet_name:
                    return pd.read_excel(file_path, sheet_name=sheet_name)
                else:
                    return pd.read_excel(file_path)
        except Exception as e:
            logger.error(f"Error reading {file_path}: {str(e)}")
            raise
    
    def update_title_slide(self, month: str, year: str):
        """Update Slide 1 title with new month/year"""
        if len(self.prs.slides) > 0:
            slide = self.prs.slides[0]
            # Find title placeholder
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text
                    if 'AIL LT' in text or 'Sep|25' in text:
                        # Replace month/year
                        month_abbr = month[:3]
                        new_text = re.sub(r'[A-Za-z]{3}\|?\d{2}', f"{month_abbr}|{year[-2:]}", text)
                        shape.text_frame.text = new_text
                        logger.info(f"Updated title slide: {new_text[:50]}")
                        break
    
    def create_fmv_slide(self, working_file_path: Path):
        """Create Slide 3 - Project FMV with exact template structure"""
        try:
            # Read CLT sheet
            df = self.read_excel_file(working_file_path, sheet_name='CLT')
            
            # Find data starting from "Division" header
            header_row = None
            for idx, row in df.iterrows():
                row_values = [str(val).strip() if pd.notna(val) else '' for val in row.values]
                if 'Division' in ' '.join(row_values):
                    header_row = idx
                    break
            
            if header_row is None:
                logger.warning("Could not find Division header in CLT sheet")
                return
            
            # Extract data - get first 10 rows of data
            data_start = header_row + 1
            table_df = df.iloc[data_start:data_start+10].copy()
            
            # Get the two columns
            if len(table_df.columns) >= 2:
                col1 = table_df.columns[0]
                col2 = table_df.columns[1]
                table_df = table_df[[col1, col2]].copy()
                table_df = table_df.dropna(subset=[col1, col2], how='all')
                table_df = table_df[table_df[col1].astype(str).str.strip() != 'Division']
                table_df = table_df[table_df[col1].astype(str).str.strip() != '']
            else:
                logger.warning("Could not find 2 columns for FMV table")
                return
            
            # Update Slide 3 (index 2)
            if len(self.prs.slides) > 2:
                slide = self.prs.slides[2]
                
                # Find the table
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table = shape.table
                        
                        # Template has 5 columns: Division Name, # of Speakers, Actual CV/FMV Status, Pending CV/FMV Status, % Response updated
                        # Excel has 2 columns: Division, Total Dis (percentage)
                        # Map: Excel col0 -> Template col0 (Division Name), Excel col1 -> Template col4 (% Response updated)
                        
                        # Fill data rows (10 data rows + 1 header = 11 total)
                        for row_idx in range(1, min(len(table.rows), len(table_df) + 1)):
                            if row_idx <= len(table_df):
                                row_data = table_df.iloc[row_idx - 1]
                                
                                # Column 0: Division Name
                                cell = table.cell(row_idx, 0)
                                div_name = str(row_data.iloc[0]) if pd.notna(row_data.iloc[0]) else ""
                                cell.text = div_name
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.LEFT
                                
                                # Column 4: % Response updated (from Excel column 1)
                                if len(row_data) > 1:
                                    value = row_data.iloc[1]
                                    if pd.notna(value):
                                        cell = table.cell(row_idx, 4)
                                        if isinstance(value, (int, float)):
                                            cell.text = f"{value:.2f}%"
                                        else:
                                            cell.text = str(value)
                                        para = cell.text_frame.paragraphs[0]
                                        para.font.size = Pt(10)
                                        para.alignment = PP_ALIGN.RIGHT
                                
                                # Leave columns 1, 2, 3 empty (not in Excel data)
                        
                        logger.info(f"Updated FMV table with {len(table_df)} rows")
                        break
        except Exception as e:
            logger.error(f"Error updating FMV slide: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def create_consent_slide(self, working_file_path: Path):
        """Create Slide 4 - Consent Status with exact template structure"""
        try:
            df = self.read_excel_file(working_file_path, sheet_name='consent')
            
            # Extract relevant columns
            relevant_cols = ['Division Name', 'DVL', '# HCP Consent', 'Consent Require', '% Consent Require']
            available_cols = [col for col in relevant_cols if col in df.columns]
            
            if len(available_cols) < 3:
                logger.warning(f"Could not find required columns in consent sheet")
                return
            
            table_df = df[available_cols].head(10).copy()
            
            # Update Slide 4 (index 3)
            if len(self.prs.slides) > 3:
                slide = self.prs.slides[3]
                
                # Find the table
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table = shape.table
                        
                        # Template has 5 columns matching our data
                        # Fill data rows
                        for row_idx in range(1, min(len(table.rows), len(table_df) + 1)):
                            if row_idx <= len(table_df):
                                row_data = table_df.iloc[row_idx - 1]
                                
                                for col_idx, col_name in enumerate(available_cols):
                                    if col_idx < len(table.columns):
                                        cell = table.cell(row_idx, col_idx)
                                        value = row_data[col_name]
                                        
                                        if pd.notna(value):
                                            if isinstance(value, float):
                                                if '%' in col_name:
                                                    cell.text = f"{value:.2f}%"
                                                else:
                                                    cell.text = f"{value:,.0f}"
                                            else:
                                                cell.text = str(value)
                                        else:
                                            cell.text = ""
                                        
                                        # Format cell
                                        para = cell.text_frame.paragraphs[0]
                                        para.font.size = Pt(10)
                                        if isinstance(value, (int, float)):
                                            para.alignment = PP_ALIGN.RIGHT
                                        else:
                                            para.alignment = PP_ALIGN.LEFT
                        
                        logger.info(f"Updated Consent table with {len(table_df)} rows")
                        break
        except Exception as e:
            logger.error(f"Error updating Consent slide: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def create_hcp_overlap_slide(self, overlap_file_path: Path):
        """Create Slide 6 - HCP Overlap with exact template structure"""
        try:
            df = self.read_excel_file(overlap_file_path)
            
            # Group by Division and count
            if 'User: Division Name' in df.columns:
                overlap_summary = df.groupby('User: Division Name').size().reset_index(name='Count')
                overlap_summary = overlap_summary.sort_values('Count', ascending=False).head(13)
            else:
                logger.warning("Could not find 'User: Division Name' column")
                return
            
            # Update Slide 6 (index 5)
            if len(self.prs.slides) > 5:
                slide = self.prs.slides[5]
                
                # Find the table
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table = shape.table
                        
                        # Template has 8 columns, we'll populate Division Name and Count
                        # Fill data rows
                        for row_idx in range(1, min(len(table.rows), len(overlap_summary) + 1)):
                            if row_idx <= len(overlap_summary):
                                row_data = overlap_summary.iloc[row_idx - 1]
                                
                                # Column 1: Division Name
                                cell = table.cell(row_idx, 1)
                                cell.text = str(row_data['User: Division Name'])
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.LEFT
                                
                                # Column 7: Count (last column)
                                if len(table.columns) > 7:
                                    cell = table.cell(row_idx, 7)
                                    cell.text = str(row_data['Count'])
                                    para = cell.text_frame.paragraphs[0]
                                    para.font.size = Pt(10)
                                    para.alignment = PP_ALIGN.RIGHT
                        
                        logger.info(f"Updated HCP Overlap table with {len(overlap_summary)} rows")
                        break
        except Exception as e:
            logger.error(f"Error updating HCP Overlap slide: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def create_missed_hcp_slide(self, chronic_missing_file_path: Path):
        """Create Slide 7 - Missed HCP with exact template structure"""
        try:
            df = self.read_excel_file(chronic_missing_file_path, sheet_name='New Visual')
            
            if 'Divison Name' in df.columns and 'Chronically missing' in df.columns:
                table_df = df[['Divison Name', 'Chronically missing', 'Strength', '%']].head(12).copy()
            else:
                logger.warning("Could not find required columns in New Visual sheet")
                return
            
            # Update Slide 7 (index 6)
            if len(self.prs.slides) > 6:
                slide = self.prs.slides[6]
                
                # Find the table
                for shape in slide.shapes:
                    if hasattr(shape, 'has_table') and shape.has_table:
                        table = shape.table
                        
                        # Template has 4 columns
                        # Fill data rows
                        for row_idx in range(1, min(len(table.rows), len(table_df) + 1)):
                            if row_idx <= len(table_df):
                                row_data = table_df.iloc[row_idx - 1]
                                
                                # Column 0: Division
                                cell = table.cell(row_idx, 0)
                                cell.text = str(row_data['Divison Name'])
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.LEFT
                                
                                # Column 1: #HCPs Missed
                                cell = table.cell(row_idx, 1)
                                cell.text = f"{int(row_data['Chronically missing']):,}"
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.RIGHT
                                
                                # Column 2: Strength
                                cell = table.cell(row_idx, 2)
                                cell.text = f"{int(row_data['Strength']):,}"
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.RIGHT
                                
                                # Column 3: %
                                cell = table.cell(row_idx, 3)
                                cell.text = f"{row_data['%']:.2f}%"
                                para = cell.text_frame.paragraphs[0]
                                para.font.size = Pt(10)
                                para.alignment = PP_ALIGN.RIGHT
                        
                        logger.info(f"Updated Missed HCP table with {len(table_df)} rows")
                        break
        except Exception as e:
            logger.error(f"Error updating Missed HCP slide: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def update_percentages(self, month: str, year: str):
        """Update percentage text in slides based on data"""
        # This will be populated from Excel data
        pass
    
    def generate_presentation(self, excel_files_dir: Path, month: str, year: str):
        """
        Generate presentation by updating template with data.
        
        Args:
            excel_files_dir: Directory containing Excel files
            month: Month name
            year: Year
        """
        excel_files_dir = Path(excel_files_dir)
        
        # Find Excel files
        working_file = excel_files_dir / "AIL LT Working file.xlsx"
        chronic_missing_file = excel_files_dir / "Chronic Missing Report AIL - Jun to Aug.xlsx"
        overlap_file = excel_files_dir / "Overlapped Vacant deactivation - greater than 2.xlsb"
        
        logger.info("Updating template with data...")
        
        # Update title slide
        self.update_title_slide(month, year)
        logger.info("Updated Slide 1: Title")
        
        # Slide 2: Business Effectiveness - keep as is (has images/charts)
        logger.info("Slide 2: Business Effectiveness (keeping template as-is)")
        
        # Slide 3: Project FMV
        if working_file.exists():
            self.create_fmv_slide(working_file)
            logger.info("Updated Slide 3: Project FMV")
        else:
            logger.warning(f"Working file not found: {working_file}")
        
        # Slide 4: Consent Status
        if working_file.exists():
            self.create_consent_slide(working_file)
            logger.info("Updated Slide 4: Consent Status")
        else:
            logger.warning(f"Working file not found: {working_file}")
        
        # Slide 5: Input Distribution - keep as is (has chart)
        logger.info("Slide 5: Input Distribution (keeping template as-is)")
        
        # Slide 6: HCP Overlap
        if overlap_file.exists():
            self.create_hcp_overlap_slide(overlap_file)
            logger.info("Updated Slide 6: HCP Overlap")
        else:
            logger.warning(f"Overlap file not found: {overlap_file}")
        
        # Slide 7: Missed HCP
        if chronic_missing_file.exists():
            self.create_missed_hcp_slide(chronic_missing_file)
            logger.info("Updated Slide 7: Missed HCP")
        else:
            logger.warning(f"Chronic missing file not found: {chronic_missing_file}")
        
        # Slide 8: Overcalled HCP - keep as is (has chart)
        logger.info("Slide 8: Overcalled HCP (keeping template as-is)")
        
        # Slide 9: Closing - keep as is
        logger.info("Slide 9: Closing (keeping template as-is)")
    
    def save(self, output_path: Path):
        """Save the presentation"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")


def get_month_year_from_input(month_input: str) -> tuple:
    """Parse month input and return month name and year"""
    month_input = month_input.strip()
    
    month_abbr_to_full = {
        'Jan': 'January', 'Feb': 'February', 'Mar': 'March', 'Apr': 'April',
        'May': 'May', 'Jun': 'June', 'Jul': 'July', 'Aug': 'August',
        'Sep': 'September', 'Oct': 'October', 'Nov': 'November', 'Dec': 'December'
    }
    
    if "'" in month_input:
        parts = month_input.split("'")
        month_abbr = parts[0].strip()
        year = "20" + parts[1].strip() if len(parts[1].strip()) == 2 else parts[1].strip()
        month_name = month_abbr_to_full.get(month_abbr, month_abbr)
        return month_name, year
    
    parts = month_input.split()
    if len(parts) >= 2:
        month_part = parts[0]
        year_part = parts[1]
        month_name = month_abbr_to_full.get(month_part, month_part)
        return month_name, year_part
    
    now = datetime.now()
    return calendar.month_name[now.month], str(now.year)


def main():
    """Main function"""
    parser = argparse.ArgumentParser(
        description='Generate AIL LT PowerPoint by replicating template exactly'
    )
    parser.add_argument(
        '--input-dir',
        type=str,
        default='excel_files',
        help='Directory containing Excel files (default: excel_files)'
    )
    parser.add_argument(
        '--output-dir',
        type=str,
        default='output',
        help='Output directory for PPT file (default: output)'
    )
    parser.add_argument(
        '--month',
        type=str,
        default=None,
        help='Month (e.g., "Sep\'25", "September 2025"). Defaults to current month'
    )
    parser.add_argument(
        '--template',
        type=str,
        default="AIL LT - Sep'25.pptx",
        help='Path to template PPT file (default: AIL LT - Sep\'25.pptx)'
    )
    parser.add_argument(
        '--output-name',
        type=str,
        default=None,
        help='Custom output filename (optional, defaults to AIL LT - [Month]\'[Year].pptx)'
    )
    
    args = parser.parse_args()
    
    # Parse month
    if args.month:
        month, year = get_month_year_from_input(args.month)
    else:
        now = datetime.now()
        month = calendar.month_name[now.month]
        year = str(now.year)
    
    logger.info("=" * 60)
    logger.info("AIL LT Template Replicator")
    logger.info("=" * 60)
    logger.info(f"Month: {month} {year}")
    logger.info(f"Template: {args.template}")
    logger.info(f"Input directory: {args.input_dir}")
    logger.info(f"Output directory: {args.output_dir}")
    
    # Initialize replicator
    template_path = Path(args.template)
    replicator = TemplateReplicator(template_path)
    
    # Generate presentation
    excel_dir = Path(args.input_dir)
    replicator.generate_presentation(excel_dir, month, year)
    
    # Save presentation
    output_dir = Path(args.output_dir)
    if args.output_name:
        output_filename = args.output_name
    else:
        month_abbr = month[:3]
        year_short = year[-2:]
        output_filename = f"AIL LT - {month_abbr}'{year_short}.pptx"
    output_path = output_dir / output_filename
    
    # Prevent overwriting template
    template_path = Path(args.template)
    if output_path.resolve() == template_path.resolve():
        # Add suffix to prevent overwriting
        output_path = output_dir / f"{output_filename.replace('.pptx', '')}_GENERATED.pptx"
        logger.warning(f"Output path matches template. Saving to: {output_path}")
    
    replicator.save(output_path)
    
    logger.info("=" * 60)
    logger.info(f"SUCCESS! Presentation saved to: {output_path}")
    logger.info(f"Total slides: {len(replicator.prs.slides)}")
    logger.info("=" * 60)
    logger.info("\nNote: The template structure, formatting, and positioning are preserved.")
    logger.info("Charts and images from the template are maintained.")


if __name__ == "__main__":
    main()

