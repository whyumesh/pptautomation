import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------- Load Excel ----------------
df = pd.read_excel("input.xlsx")

# ---------------- Create Chart using Seaborn ----------------
plt.figure(figsize=(8, 4))
sns.barplot(data=df, x="Metric", y="Value", palette="deep")
plt.title("Key Metrics Overview")
plt.xticks(rotation=20)
plt.tight_layout()

plt.savefig("chart.png")
plt.close()

# ---------------- Create PPT ----------------
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Monthly Business Report"
slide.placeholders[1].text = "Automated Excel → Visualization → PPT"

# Data Summary Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Metric Summary"

tf = slide.shapes.placeholders[1].text_frame
tf.clear()

for _, row in df.iterrows():
    p = tf.add_paragraph()
    p.text = f"{row['Metric']}: {row['Value']}"
    p.font.size = Pt(18)

# Chart Slide
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
slide.shapes.title.text = "Visual Insights"

slide.shapes.add_picture(
    "chart.png",
    Inches(1),
    Inches(1.5),
    width=Inches(8)
)

# Conclusion Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
slide.placeholders[1].text = (
    "• Scale to multiple Excel files\n"
    "• Map charts to fixed PPT template\n"
    "• Automate monthly execution"
)

prs.save("output.pptx")

print("✅ PPT with Seaborn chart generated successfully")
